from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_ai_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Set dark theme for all slides
    def set_slide_background(slide, color_rgb=(26, 35, 126)):  # Dark navy blue
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color_rgb)
    
    def add_diagram_shape(slide, shape_type, left, top, width, height, text="", color_rgb=(255, 255, 255)):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color_rgb)
        if text:
            shape.text = text
            shape.text_frame.paragraphs[0].font.size = Pt(12)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        return shape
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Artificial Intelligence in Crime Pattern Prediction"
    subtitle.text = "A Machine Learning Approach for Tamil Nadu\nPresented by: [Your Name]\nDate: [Current Date]"
    
    # Set dark background
    set_slide_background(slide)
    
    # Make text white
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Slide 2: What is Artificial Intelligence?
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What is Artificial Intelligence?"
    content.text = """• AI is the simulation of human intelligence in machines
• Enables computers to learn, reason, and make decisions
• Key components:
  - Machine Learning (ML)
  - Deep Learning
  - Natural Language Processing
  - Computer Vision
• Applications: Healthcare, Finance, Transportation, Security"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add AI components diagram
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.5), Inches(3), Inches(0.8), "Machine Learning", (100, 149, 237))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(2.5), Inches(3), Inches(0.8), "Deep Learning", (72, 209, 204))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(3.5), Inches(3), Inches(0.8), "NLP", (255, 182, 193))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(4.5), Inches(3), Inches(0.8), "Computer Vision", (255, 218, 185))
    
    # Slide 3: Machine Learning Overview
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Machine Learning in Crime Prediction"
    content.text = """• Supervised Learning: Learning from labeled data
• Classification: Predicting categories (crime types)
• Algorithms Used:
  - Logistic Regression
  - Decision Trees
  - Random Forest
• Benefits:
  - Pattern Recognition
  - Predictive Analytics
  - Resource Optimization"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add ML workflow diagram
    add_diagram_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(1), Inches(2.5), Inches(0.8), "Data Input", (144, 238, 144))
    # Arrow
    add_diagram_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(7), Inches(2), Inches(1), Inches(0.5), "", (255, 255, 255))
    add_diagram_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(2.8), Inches(2.5), Inches(0.8), "ML Algorithm", (255, 165, 0))
    # Arrow
    add_diagram_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(7), Inches(3.8), Inches(1), Inches(0.5), "", (255, 255, 255))
    add_diagram_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(4.6), Inches(2.5), Inches(0.8), "Prediction", (255, 99, 132))
    
    # Slide 4: Project Overview
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Tamil Nadu Crime Pattern Predictor"
    content.text = """• Objective: Predict crime types based on various factors
• Dataset: 7,000 synthetic crime records
• Features: 25 attributes including:
  - Geographic data (District, Taluk)
  - Temporal data (Time, Day, Month)
  - Demographic data (Age Group)
  - Environmental factors (Area Type, Public Events)
• Target: 10 different crime types"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add Tamil Nadu map representation
    add_diagram_shape(slide, MSO_SHAPE.HEXAGON, Inches(6.5), Inches(1.5), Inches(2), Inches(1.5), "Tamil Nadu\n22 Districts", (70, 130, 180))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5), Inches(1.5), Inches(0.6), "7K Records", (255, 215, 0))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(3.5), Inches(1.5), Inches(0.6), "25 Features", (255, 215, 0))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(4.5), Inches(1.5), Inches(0.6), "10 Crime Types", (255, 99, 132))
    
    # Slide 5: Dataset Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Dataset Features & Structure"
    content.text = """Geographic Features:
• District (22 districts in Tamil Nadu)
• Taluk (Sub-districts)
• Area Type (Urban/Semi-Urban/Rural)

Temporal Features:
• Time of Day (Morning/Afternoon/Evening/Night)
• Day of Week (Monday-Sunday)
• Month (1-12)

Demographic & Social:
• Age Group (18-25, 26-35, 36-45, 46-60, 60+)
• Public Event (0/1)
• Population Density, Literacy Rate"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add feature categories diagram
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(1), Inches(2.5), Inches(1), "Geographic\nFeatures", (135, 206, 235))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(2.5), Inches(2.5), Inches(1), "Temporal\nFeatures", (255, 182, 193))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(4), Inches(2.5), Inches(1), "Demographic\nFeatures", (144, 238, 144))
    
    # Slide 6: ML Pipeline
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Machine Learning Pipeline"
    content.text = """1. Data Generation
   • Synthetic dataset creation with realistic patterns
   
2. Data Preprocessing
   • Feature encoding (Label/One-hot encoding)
   • Data scaling and normalization
   
3. Model Training
   • Multiple algorithm comparison
   • Hyperparameter optimization
   
4. Model Evaluation
   • Accuracy metrics
   • Confusion matrix analysis
   
5. Deployment
   • Flask web application
   • Real-time predictions"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add pipeline flow diagram
    steps = [("Data Gen", (144, 238, 144)), ("Preprocess", (255, 182, 193)), 
             ("Training", (255, 165, 0)), ("Evaluation", (135, 206, 235)), ("Deploy", (255, 99, 132))]
    
    for i, (step, color) in enumerate(steps):
        add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6 + i*0.4), Inches(1.5 + i*0.6), 
                         Inches(1.8), Inches(0.6), step, color)
        if i < len(steps) - 1:
            add_diagram_shape(slide, MSO_SHAPE.RIGHT_ARROW, Inches(6.2 + i*0.4), Inches(2.3 + i*0.6), 
                             Inches(0.6), Inches(0.3), "", (255, 255, 255))
    
    # Slide 7: Model Performance
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Model Performance Results"
    content.text = """Algorithm Comparison:
• Decision Tree: 100% accuracy (100% training data)
• Random Forest: 95%+ accuracy
• Logistic Regression: ~18% accuracy (with train/test split)

Key Insights:
• Location (District/Taluk) is most predictive feature
• Time patterns significantly influence crime types
• Age group and public events are important factors
• Overfitting observed with 100% training approach"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add performance bar chart representation
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.5), Inches(0.8), Inches(2.5), "Decision\nTree\n100%", (50, 205, 50))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(0.8), Inches(2), "Random\nForest\n95%", (255, 165, 0))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8), Inches(3.5), Inches(0.8), Inches(0.5), "Logistic\nRegression\n18%", (255, 99, 132))
    
    # Slide 8: Web Application
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Web Application Features"
    content.text = """User Interface:
• Clean, professional design with Bootstrap
• Responsive layout for all devices

Prediction Form:
• Input fields for all relevant features
• Dynamic taluk selection based on district
• Real-time crime type prediction

Analytics Dashboard:
• Crime distribution visualizations
• District-wise analysis charts
• Time pattern analysis
• Interactive data exploration"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add web app architecture diagram
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(1), Inches(2.5), Inches(0.8), "Frontend (HTML/CSS/JS)", (135, 206, 235))
    add_diagram_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(7), Inches(2), Inches(1), Inches(0.4), "", (255, 255, 255))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(2.8), Inches(2.5), Inches(0.8), "Flask Backend", (255, 165, 0))
    add_diagram_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(7), Inches(3.8), Inches(1), Inches(0.4), "", (255, 255, 255))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(4.6), Inches(2.5), Inches(0.8), "ML Model", (255, 99, 132))
    
    # Slide 9: Technical Implementation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Implementation"
    content.text = """Backend Technologies:
• Python Flask framework
• Scikit-learn for ML algorithms
• Pandas for data manipulation
• Joblib for model serialization

Frontend Technologies:
• HTML5, CSS3, Bootstrap 4
• JavaScript for dynamic interactions
• Chart.js for data visualizations

Data Processing:
• NumPy for numerical computations
• Matplotlib/Seaborn for plotting
• Custom feature engineering pipeline"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add tech stack diagram
    add_diagram_shape(slide, MSO_SHAPE.HEXAGON, Inches(6), Inches(1), Inches(1.2), Inches(0.8), "Python", (255, 215, 0))
    add_diagram_shape(slide, MSO_SHAPE.HEXAGON, Inches(7.5), Inches(1), Inches(1.2), Inches(0.8), "Flask", (144, 238, 144))
    add_diagram_shape(slide, MSO_SHAPE.HEXAGON, Inches(6), Inches(2.2), Inches(1.2), Inches(0.8), "Scikit-learn", (255, 182, 193))
    add_diagram_shape(slide, MSO_SHAPE.HEXAGON, Inches(7.5), Inches(2.2), Inches(1.2), Inches(0.8), "Bootstrap", (135, 206, 235))
    add_diagram_shape(slide, MSO_SHAPE.HEXAGON, Inches(6), Inches(3.4), Inches(1.2), Inches(0.8), "Pandas", (255, 165, 0))
    add_diagram_shape(slide, MSO_SHAPE.HEXAGON, Inches(7.5), Inches(3.4), Inches(1.2), Inches(0.8), "Chart.js", (255, 99, 132))
    
    # Slide 10: AI Prompts Used
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "AI Development Prompts"
    content.text = """Data Generation Prompt:
"Create a synthetic crime dataset for Tamil Nadu with 7,000 records including geographic, temporal, and demographic features"

Model Selection Prompt:
"Compare multiple ML algorithms for crime classification and recommend the best approach"

Feature Engineering Prompt:
"Design meaningful features from raw crime data that can improve prediction accuracy"

UI/UX Design Prompt:
"Create a clean, professional web interface for crime prediction with user-friendly forms and visualizations" """
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add AI prompt workflow
    add_diagram_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(1), Inches(2.5), Inches(0.6), "Data Prompt", (144, 238, 144))
    add_diagram_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(2), Inches(2.5), Inches(0.6), "Model Prompt", (255, 182, 193))
    add_diagram_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(3), Inches(2.5), Inches(0.6), "Feature Prompt", (255, 165, 0))
    add_diagram_shape(slide, MSO_SHAPE.OVAL, Inches(6), Inches(4), Inches(2.5), Inches(0.6), "UI/UX Prompt", (135, 206, 235))
    
    # Slide 11: Future Enhancements
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Enhancements"
    content.text = """Data Improvements:
• Integration with real crime databases
• Larger dataset for better generalization
• Real-time data feeds

Model Enhancements:
• Deep learning approaches
• Ensemble methods
• Time series forecasting

Deployment:
• Cloud platform deployment (AWS/Azure)
• Mobile application development
• API for third-party integration
• Real-time alert system"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Add future roadmap diagram
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.5), Inches(2.5), Inches(0.8), "Real Data Integration", (144, 238, 144))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(2.8), Inches(2.5), Inches(0.8), "Deep Learning Models", (255, 182, 193))
    add_diagram_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6), Inches(4.1), Inches(2.5), Inches(0.8), "Cloud Deployment", (135, 206, 235))
    
    # Slide 12: Sources & References
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Sources & References"
    content.text = """Academic Sources:
• Scikit-learn Documentation: https://scikit-learn.org/
• "Pattern Recognition and Machine Learning" - Bishop, C.M.
• "Hands-On Machine Learning" - Aurélien Géron

Technical Resources:
• Flask Documentation: https://flask.palletsprojects.com/
• Bootstrap Framework: https://getbootstrap.com/
• Chart.js Documentation: https://www.chartjs.org/

Crime Data References:
• National Crime Records Bureau (NCRB) India
• Tamil Nadu Police Department Statistics
• Crime Pattern Analysis Research Papers"""
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Slide 13: Thank You
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You"
    subtitle.text = """Questions & Discussion
    
Contact Information:
Email: [your.email@domain.com]
GitHub: [your-github-profile]
Project Repository: [repository-link]

"AI is not about replacing human intelligence,
but augmenting it for better decision making" """
    
    # Set dark background and white text
    set_slide_background(slide)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Save presentation
    prs.save("AI_Crime_Prediction_Dark_Theme.pptx")
    print("Presentation saved as: AI_Crime_Prediction_Dark_Theme.pptx")

if __name__ == "__main__":
    create_ai_presentation()